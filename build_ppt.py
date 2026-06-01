"""Generate SpeakUp capstone tech-stack PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ----
NAVY = RGBColor(0x0F, 0x1F, 0x3A)
ACCENT = RGBColor(0x4F, 0x8C, 0xFF)
ACCENT_SOFT = RGBColor(0xE8, 0xF0, 0xFE)
GREY_DARK = RGBColor(0x2B, 0x33, 0x40)
GREY_MID = RGBColor(0x5A, 0x63, 0x70)
GREY_LIGHT = RGBColor(0xEE, 0xF1, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xFF, 0xC4, 0x4D)
RED = RGBColor(0xE5, 0x4B, 0x4B)

FONT_KR = "Malgun Gothic"
FONT_MONO = "Consolas"

# ---- layout ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill, line=None, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=GREY_DARK,
             font=FONT_KR, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rich(slide, x, y, w, h, runs, size=14, color=GREY_DARK, font=FONT_KR,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs: list of (text, {bold, color, size, font, italic})"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, opts in runs:
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            continue
        run = p.add_run()
        run.text = text
        run.font.name = opts.get("font", font)
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", color)
    return tb


def add_bullets(slide, x, y, w, h, bullets, size=15, color=GREY_DARK,
                bullet_color=ACCENT, line_spacing=1.3, bold_first=False):
    """bullets: list of strings (use '**text**' for bold sub-spans)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet dot
        dot = p.add_run()
        dot.text = "•  "
        dot.font.name = FONT_KR
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = bullet_color
        # parse **bold** segments
        segments = []
        cur = ""
        j = 0
        bold = False
        while j < len(line):
            if line[j:j+2] == "**":
                if cur:
                    segments.append((cur, bold))
                    cur = ""
                bold = not bold
                j += 2
            else:
                cur += line[j]
                j += 1
        if cur:
            segments.append((cur, bold))
        for seg, b in segments:
            r = p.add_run()
            r.text = seg
            r.font.name = FONT_KR
            r.font.size = Pt(size)
            r.font.bold = b or (bold_first and segments.index((seg, b)) == 0)
            r.font.color.rgb = NAVY if b else color
    return tb


def header_bar(slide, title, subtitle=None, kicker=None):
    # left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.22), SLIDE_H, ACCENT)
    # top kicker
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.35), Inches(8), Inches(0.3),
                 kicker, size=11, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.8),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.25), Inches(12), Inches(0.5),
                 subtitle, size=14, color=GREY_MID)
    # divider
    add_rect(slide, Inches(0.6), Inches(1.78), Inches(12.1), Emu(15000), GREY_LIGHT)


def footer(slide, page, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "SpeakUp — 멀티모달 커뮤니케이션 코칭 플랫폼", size=9, color=GREY_MID)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{page} / {total}", size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)


TOTAL = 14

# =========================================================================
# Slide 1 — Title
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
# decorative circles
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(-1.5), Inches(5), Inches(5))
c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
c.fill.transparency = 0.5
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5), Inches(4), Inches(4))
c2.fill.solid(); c2.fill.fore_color.rgb = ACCENT; c2.line.fill.background()

add_text(s, Inches(0.9), Inches(1.2), Inches(8), Inches(0.4),
         "CAPSTONE 2026 · TECH DEEP-DIVE", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(1.7), Inches(12), Inches(1.6),
         "SpeakUp", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(3.4), Inches(12), Inches(1.2),
         "멀티모달 커뮤니케이션 코칭 플랫폼", size=30, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(4.3), Inches(11), Inches(0.6),
         "발표 · 면접 · 보컬 — 같은 엔진, 루브릭만 교체", size=18, color=ACCENT_SOFT)
# tech chips
chips = ["MediaPipe", "Web Speech API", "faster-whisper", "@pixiv/three-vrm",
         "PostgreSQL", "FastAPI", "IndexedDB"]
cx = Inches(0.9)
cy = Inches(5.4)
for ch in chips:
    w = Inches(0.18 + 0.13 * len(ch))
    add_round(s, cx, cy, w, Inches(0.42), ACCENT_SOFT, radius=0.5)
    add_text(s, cx, cy, w, Inches(0.42), ch, size=11, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w + Inches(0.12)

add_text(s, Inches(0.9), Inches(6.6), Inches(8), Inches(0.4),
         "팀 Sprinter · 2026", size=14, color=ACCENT_SOFT)

# =========================================================================
# Slide 2 — 문제 & 차별점
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "기존 코칭 앱과 무엇이 다른가", kicker="PROBLEM & DIFFERENTIATION")

# left: 기존
add_round(s, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.0), GREY_LIGHT)
add_text(s, Inches(0.9), Inches(2.25), Inches(5.4), Inches(0.4),
         "기존 코칭 앱", size=13, bold=True, color=GREY_MID)
add_text(s, Inches(0.9), Inches(2.7), Inches(5.4), Inches(1.4),
         "영상 업로드 → 사후 한 줄 평가\n신호는 단일 채널 (대개 음성만)\n시나리오 추가 = 새 앱",
         size=15, color=GREY_DARK)

# right: ours
add_round(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.0), ACCENT_SOFT)
add_text(s, Inches(7.1), Inches(2.25), Inches(5.4), Inches(0.4),
         "SpeakUp", size=13, bold=True, color=ACCENT)
add_text(s, Inches(7.1), Inches(2.7), Inches(5.4), Inches(1.4),
         "발화 한 문장 끝나는 순간 LLM 코칭\n시선·표정·자세·손·운율·텍스트 동시\nYAML 한 장 추가 = 새 시나리오",
         size=15, bold=True, color=NAVY)

# 3 cards
cards = [
    ("즉시 반응", "Web Speech API isFinal 단위 LLM 호출", ACCENT),
    ("다중 신호", "비전 5fps + 운율 1Hz + 텍스트", HIGHLIGHT),
    ("확장성", "rubrics/*.yaml 만 교체", RED),
]
x0 = Inches(0.6)
for i, (t, d, col) in enumerate(cards):
    cx = x0 + Inches(i * 4.1)
    add_round(s, cx, Inches(4.5), Inches(3.9), Inches(2.0), WHITE, line=GREY_LIGHT)
    add_rect(s, cx, Inches(4.5), Inches(3.9), Inches(0.18), col)
    add_text(s, cx + Inches(0.3), Inches(4.85), Inches(3.5), Inches(0.5),
             t, size=20, bold=True, color=NAVY)
    add_text(s, cx + Inches(0.3), Inches(5.5), Inches(3.5), Inches(1.0),
             d, size=13, color=GREY_MID)

footer(s, 2, TOTAL)

# =========================================================================
# Slide 3 — System Overview L1/L2/L3
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "3-Layer 평가 구조", subtitle="실시간 HUD + 의미 이벤트 + 사후 종합 리포트가 한 세션에서 동시에",
           kicker="SYSTEM OVERVIEW")

layers = [
    ("L1", "Rule Engine", "즉시 HUD", "7개 트리거 (silence / gaze / smile_absence /\nmotion_absence / filler / speech_rate / content)\n각자 cooldown · 글로벌 2s gate", ACCENT),
    ("L2", "Semantic Events", "의미 단위", "17종 이벤트로 정규화\nchin_on_hand · head_tilt_sustained · …\nVisionContextBuffer 3s 스무딩", HIGHLIGHT),
    ("L3", "LLM Report", "사후 종합", "시나리오 YAML 루브릭 로드\naxis 점수 + annotated moments\n+ 훈련 처방 (Gemini / Claude)", RED),
]
x0 = Inches(0.6)
for i, (tag, name, sub, body, col) in enumerate(layers):
    cx = x0 + Inches(i * 4.1)
    add_round(s, cx, Inches(2.2), Inches(3.9), Inches(4.3), WHITE, line=GREY_LIGHT)
    add_rect(s, cx, Inches(2.2), Inches(3.9), Inches(0.5), col)
    add_text(s, cx + Inches(0.3), Inches(2.27), Inches(3.5), Inches(0.4),
             tag, size=18, bold=True, color=WHITE)
    add_text(s, cx + Inches(0.3), Inches(2.95), Inches(3.5), Inches(0.5),
             name, size=22, bold=True, color=NAVY)
    add_text(s, cx + Inches(0.3), Inches(3.55), Inches(3.5), Inches(0.4),
             sub, size=12, bold=True, color=col)
    add_text(s, cx + Inches(0.3), Inches(4.05), Inches(3.5), Inches(2.3),
             body, size=12, color=GREY_DARK)
    # arrow
    if i < 2:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 cx + Inches(3.95), Inches(4.25), Inches(0.18), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = GREY_MID; arr.line.fill.background()

add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3),
         "왼쪽 = 실시간   ·   오른쪽 = 사후   —   둘 다 같은 세션에서 병렬로 흐릅니다.",
         size=12, color=GREY_MID, align=PP_ALIGN.CENTER)
footer(s, 3, TOTAL)

# =========================================================================
# Slide 4 — Tech Stack Map (table)
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "Tech Stack — 전 영역 오픈소스 조합", kicker="STACK MAP")

rows = [
    ("비전 (브라우저)", "MediaPipe Tasks-Vision (WASM, GPU)", "Face·Pose·Hand Landmarker"),
    ("실시간 STT", "Web Speech API", "브라우저 네이티브 · isFinal 단위"),
    ("사후 STT", "faster-whisper (medium, int8)", "word-timestamp + VAD"),
    ("운율 분석", "librosa", "F0 / RMS / pitch SD"),
    ("아바타", "@pixiv/three-vrm + three.js", "VRM 로더·뼈대 API"),
    ("녹화", "MediaRecorder API", "video/webm; vp9+opus"),
    ("클라이언트 저장", "IndexedDB", "대용량 blob 영구 보관"),
    ("서버 저장", "PostgreSQL 16 (JSONB)", "메타·세션·LLM 메시지"),
    ("백엔드", "FastAPI × 3 + Docker Compose", "audio · aggregator · coach"),
    ("LLM", "Gemini 2.5 Flash · Claude Sonnet 4.6", "LLM_PROVIDER로 스왑"),
]
y = Inches(2.0)
row_h = Inches(0.45)
col_w = [Inches(2.6), Inches(4.6), Inches(4.9)]
xs = [Inches(0.6), Inches(0.6) + col_w[0], Inches(0.6) + col_w[0] + col_w[1]]

# header row
add_rect(s, xs[0], y, col_w[0] + col_w[1] + col_w[2], row_h, NAVY)
headers = ["레이어", "기술", "역할 / 비고"]
for i, h in enumerate(headers):
    add_text(s, xs[i] + Inches(0.15), y, col_w[i], row_h,
             h, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
y += row_h

for idx, (a, b, c) in enumerate(rows):
    bg = GREY_LIGHT if idx % 2 == 0 else WHITE
    add_rect(s, xs[0], y, col_w[0] + col_w[1] + col_w[2], row_h, bg)
    add_text(s, xs[0] + Inches(0.15), y, col_w[0], row_h,
             a, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, xs[1] + Inches(0.15), y, col_w[1], row_h,
             b, size=11, color=GREY_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, xs[2] + Inches(0.15), y, col_w[2], row_h,
             c, size=11, color=GREY_MID, anchor=MSO_ANCHOR.MIDDLE)
    y += row_h

add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
         "자체 모델 학습 0 · 핵심 무기는 '잘 고른 조합'과 그 위의 시나리오 추상화",
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 4, TOTAL)

# =========================================================================
# Slide 5 — 실시간 STT & 에이전트 리액션 ★
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "실시간 STT → 에이전트 리액션",
           subtitle="Web Speech API의 isFinal = 문장 경계. 문장이 끝나는 그 순간에만 LLM이 깨어납니다.",
           kicker="★ CORE MECHANISM")

# left bullets
add_bullets(s, Inches(0.6), Inches(2.1), Inches(6.3), Inches(4.5), [
    "**Web Speech API**의 `isFinal` 플래그를 발화 경계로 사용",
    "단어가 아닌 **'한 문장'** 단위로 LLM 1회 호출",
    "프롬프트는 두 부분으로 구성:",
    "    └  **이전 문장들** = context (누적 transcript)",
    "    └  **현재 문장** = 평가/리액션 대상",
    "발화 종료 직후 **≈ 1초 내** 코치 코멘트",
    "글로벌 **2s cooldown** + selfHistory 안티-반복으로 비용 제어",
], size=14, line_spacing=1.45)

# right diagram
dx = Inches(7.3)
dy = Inches(2.1)
dw = Inches(5.4)

# context box
add_round(s, dx, dy, dw, Inches(1.4), GREY_LIGHT)
add_text(s, dx + Inches(0.25), dy + Inches(0.1), dw, Inches(0.35),
         "context  ·  이전 문장들", size=10, bold=True, color=GREY_MID)
add_text(s, dx + Inches(0.25), dy + Inches(0.5), dw - Inches(0.5), Inches(0.9),
         "“여러분 안녕하세요. 오늘 발표는 …”\n“먼저 배경부터 말씀드리면 …”",
         size=12, color=GREY_DARK)

# arrow down
a1 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                        dx + Inches(2.55), dy + Inches(1.45),
                        Inches(0.3), Inches(0.25))
a1.fill.solid(); a1.fill.fore_color.rgb = ACCENT; a1.line.fill.background()

# current box
add_round(s, dx, dy + Inches(1.8), dw, Inches(1.0), ACCENT_SOFT, line=ACCENT)
add_text(s, dx + Inches(0.25), dy + Inches(1.88), dw, Inches(0.35),
         "current  ·  isFinal로 막 끝난 문장", size=10, bold=True, color=ACCENT)
add_text(s, dx + Inches(0.25), dy + Inches(2.22), dw - Inches(0.5), Inches(0.6),
         "“어… 그러니까 결론부터 말씀드리면”",
         size=14, bold=True, color=NAVY)

a2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                        dx + Inches(2.55), dy + Inches(2.85),
                        Inches(0.3), Inches(0.25))
a2.fill.solid(); a2.fill.fore_color.rgb = ACCENT; a2.line.fill.background()

# LLM
add_round(s, dx + Inches(1.4), dy + Inches(3.2), Inches(2.6), Inches(0.55),
          NAVY, radius=0.5)
add_text(s, dx + Inches(1.4), dy + Inches(3.2), Inches(2.6), Inches(0.55),
         "LLM (Gemini / Claude)", size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

a3 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                        dx + Inches(2.55), dy + Inches(3.8),
                        Inches(0.3), Inches(0.25))
a3.fill.solid(); a3.fill.fore_color.rgb = ACCENT; a3.line.fill.background()

add_round(s, dx, dy + Inches(4.15), dw, Inches(0.6), HIGHLIGHT, radius=0.4)
add_text(s, dx, dy + Inches(4.15), dw, Inches(0.6),
         "코치 코멘트 1개 ↩  (≈ 1s)", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 5, TOTAL)

# =========================================================================
# Slide 6 — 비전 파이프라인
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "비전 파이프라인 — MediaPipe Tasks-Vision",
           subtitle="브라우저 WASM에서 5 FPS 추론. 원본 프레임은 절대 서버로 가지 않습니다.",
           kicker="VISION")

# left: three landmarker cards
items = [
    ("FaceLandmarker", "478 landmarks\n52 blendshapes\nhead transform matrix",
     "시선 · 미소 · 끄덕임 · 표정"),
    ("PoseLandmarker", "33 keypoints (상체 위주)",
     "어깨 기울기 · 자세 · 턱괴기"),
    ("HandLandmarker", "21 keypoints × 2",
     "손짓 빈도 · 제스처"),
]
y = Inches(2.1)
for i, (t, d, use) in enumerate(items):
    cy = y + Inches(i * 1.5)
    add_round(s, Inches(0.6), cy, Inches(7.2), Inches(1.35), WHITE, line=GREY_LIGHT)
    add_rect(s, Inches(0.6), cy, Inches(0.16), Inches(1.35), ACCENT)
    add_text(s, Inches(0.95), cy + Inches(0.15), Inches(2.8), Inches(0.4),
             t, size=15, bold=True, color=NAVY)
    add_text(s, Inches(0.95), cy + Inches(0.55), Inches(2.8), Inches(0.8),
             d, size=11, color=GREY_MID)
    add_text(s, Inches(4.2), cy + Inches(0.15), Inches(3.5), Inches(0.4),
             "산출 신호", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(4.2), cy + Inches(0.5), Inches(3.5), Inches(0.8),
             use, size=13, bold=True, color=GREY_DARK)

# right: privacy panel
add_round(s, Inches(8.1), Inches(2.1), Inches(4.6), Inches(4.4),
          NAVY, radius=0.05)
add_text(s, Inches(8.35), Inches(2.3), Inches(4.2), Inches(0.5),
         "PRIVACY BY DESIGN", size=11, bold=True, color=ACCENT)
add_text(s, Inches(8.35), Inches(2.75), Inches(4.2), Inches(0.7),
         "원본 카메라 프레임은\n브라우저 밖으로 나가지 않음",
         size=15, bold=True, color=WHITE)

add_rect(s, Inches(8.35), Inches(4.05), Inches(4.1), Pt(1), ACCENT)

add_text(s, Inches(8.35), Inches(4.2), Inches(4.2), Inches(0.4),
         "서버로 전송되는 것은", size=11, color=ACCENT_SOFT)
add_text(s, Inches(8.35), Inches(4.55), Inches(4.2), Inches(1.8),
         "✓ 추출된 랜드마크 숫자\n✓ blendshape 강도\n✓ 트리거 이벤트 JSON\n→  ≈ 1 KB / frame",
         size=13, color=WHITE)

footer(s, 6, TOTAL)

# =========================================================================
# Slide 7 — VRM 경계 ★
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "VRM 아바타 — 라이브러리는 어디까지, 우리는 어디부터",
           subtitle="@pixiv/three-vrm 은 '아바타를 다룰 API'를 줍니다. 카메라 사람을 따라 움직이게 하는 건 우리가 짭니다.",
           kicker="★ OPEN-SOURCE BOUNDARY")

# library card (left)
add_round(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(4.5),
          ACCENT_SOFT, line=ACCENT)
add_text(s, Inches(0.85), Inches(2.25), Inches(5.5), Inches(0.4),
         "@pixiv/three-vrm  +  three.js", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(2.7), Inches(5.5), Inches(0.5),
         "라이브러리가 대신해주는 것", size=18, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(3.4), Inches(5.5), Inches(3.0),
         "✓  VRM 1.0 파일 파싱 · 렌더링\n"
         "✓  Humanoid Skeleton 추상화\n"
         "      vrm.humanoid.getNormalizedBoneNode('leftUpperArm')\n"
         "✓  BlendShape (표정) 인터페이스\n"
         "✓  LookAt (시선) 인터페이스",
         size=13, color=GREY_DARK)

# our code card (right)
add_round(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.5),
          WHITE, line=RED)
add_text(s, Inches(7.05), Inches(2.25), Inches(5.5), Inches(0.4),
         "우리가 직접 짠 부분", size=14, bold=True, color=RED)
add_text(s, Inches(7.05), Inches(2.7), Inches(5.5), Inches(0.5),
         "포즈 따라하기는 라이브러리가 안 함", size=18, bold=True, color=NAVY)
add_text(s, Inches(7.05), Inches(3.4), Inches(5.5), Inches(3.0),
         "✗  MediaPipe Pose 33 keypoints  →  VRM bone Quaternion\n"
         "        (retargeting 로직)\n"
         "✗  MediaPipe Face 52 blendshape  →  VRM expression 매핑\n"
         "✗  canvas.captureStream(30)으로 아바타 별도 녹화\n"
         "✗  좌우 반전 · 어깨 안정화 등 보정",
         size=13, color=GREY_DARK)

# bottom flow
add_text(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.3),
         "MediaPipe Landmarks   →   [retarget 자체 구현]   →   VRM Bone Rotation   →   three.js Renderer",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
footer(s, 7, TOTAL)

# =========================================================================
# Slide 8 — 실시간 측정이 어떻게 가능한가
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "실시간으로 평가 신호가 흐르는 방식",
           subtitle="원본 영상이 아니라 '압축된 의미 신호'만 흐르기 때문에 노트북 한 대로 충분",
           kicker="REAL-TIME PIPELINE")

# 5 channels on a timeline
channels = [
    ("Face", "5 FPS", ACCENT),
    ("Pose", "5 FPS", ACCENT),
    ("Hand", "5 FPS", ACCENT),
    ("Prosody", "1 Hz", HIGHLIGHT),
    ("Text", "isFinal", RED),
]
x0 = Inches(0.6)
ch_w = Inches(2.45)
for i, (n, r, col) in enumerate(channels):
    cx = x0 + Inches(i * 2.5)
    add_round(s, cx, Inches(2.2), ch_w, Inches(1.2), WHITE, line=GREY_LIGHT)
    add_rect(s, cx, Inches(2.2), Inches(0.12), Inches(1.2), col)
    add_text(s, cx + Inches(0.25), Inches(2.35), ch_w, Inches(0.4),
             n, size=15, bold=True, color=NAVY)
    add_text(s, cx + Inches(0.25), Inches(2.75), ch_w, Inches(0.3),
             r, size=11, bold=True, color=col)
    add_text(s, cx + Inches(0.25), Inches(3.0), ch_w, Inches(0.4),
             "≈ 1 KB/frame", size=10, color=GREY_MID)

# arrow → ws
add_text(s, Inches(0.6), Inches(3.65), Inches(12), Inches(0.4),
         "↓   WebSocket (JSON)", size=12, bold=True, color=GREY_MID, align=PP_ALIGN.CENTER)

# Aggregator box
add_round(s, Inches(2), Inches(4.2), Inches(9.3), Inches(1.4), NAVY, radius=0.1)
add_text(s, Inches(2), Inches(4.3), Inches(9.3), Inches(0.4),
         "Aggregator  (:8001)", size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.3), Inches(4.7), Inches(8.7), Inches(0.9),
         "VisionContextBuffer (3s 스무딩)   +   TranscriptBuffer   +   7 Trigger Evaluators   +   2s global cooldown",
         size=12, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# arrow → coach
add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
         "↓   임계 통과한 트리거만", size=12, bold=True, color=GREY_MID, align=PP_ALIGN.CENTER)

add_round(s, Inches(4.3), Inches(6.15), Inches(4.7), Inches(0.55),
          HIGHLIGHT, radius=0.4)
add_text(s, Inches(4.3), Inches(6.15), Inches(4.7), Inches(0.55),
         "Coach LLM  →  실시간 HUD 코멘트", size=13, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 8, TOTAL)

# =========================================================================
# Slide 9 — 대용량 영상 처리 (no AWS) ★
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "대용량 녹화 영상 처리 — AWS 없이 어떻게?",
           subtitle="설계 원칙: '영상은 사용자 디바이스에 산다.' 서버는 잠깐 분석만 하고 폐기.",
           kicker="★ STORAGE & UPLOAD")

# Pipeline 4 steps
steps = [
    ("1", "녹화", "MediaRecorder API\nvideo/webm; vp9+opus\n+ canvas.captureStream(30)\n(아바타 별도 트랙)"),
    ("2", "클라이언트 저장", "IndexedDB\nspeakup-media DB\n→ Blob 영구 보관\n(localStorage = 메타만)"),
    ("3", "업로드", "POST /analyze\nmultipart 단일 POST\n청크 분할 ✗\n로컬 Docker 서비스로 직행"),
    ("4", "서버 처리", "ffmpeg → 16kHz mono wav\nfaster-whisper + librosa\n분석 끝 → 원본 webm 폐기\n결과만 PostgreSQL JSONB"),
]
x0 = Inches(0.6)
step_w = Inches(2.95)
for i, (n, t, d) in enumerate(steps):
    cx = x0 + Inches(i * 3.05)
    add_round(s, cx, Inches(2.1), step_w, Inches(2.6), WHITE, line=GREY_LIGHT)
    # number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.2), Inches(2.25),
                              Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background()
    add_text(s, cx + Inches(0.2), Inches(2.25), Inches(0.5), Inches(0.5),
             n, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx + Inches(0.85), Inches(2.3), step_w, Inches(0.4),
             t, size=15, bold=True, color=NAVY)
    add_text(s, cx + Inches(0.25), Inches(2.95), step_w - Inches(0.5), Inches(1.7),
             d, size=11, color=GREY_DARK)
    if i < 3:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 cx + step_w + Inches(0.0), Inches(3.3),
                                 Inches(0.15), Inches(0.25))
        arr.fill.solid(); arr.fill.fore_color.rgb = GREY_MID; arr.line.fill.background()

# bottom: key design decisions
add_round(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(1.7),
          ACCENT_SOFT, line=ACCENT)
add_text(s, Inches(0.85), Inches(5.1), Inches(11.5), Inches(0.4),
         "핵심 설계 결정", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(5.5), Inches(11.5), Inches(1.1),
         "✓  S3 / 클라우드 스토리지 ✗   →   서버는 영상 영구 저장하지 않음\n"
         "✓  청크 분할 업로드 ✗   →   단일 multipart POST (대역폭·타임아웃이 사실상 한도, 로컬 검증 수백 MB)\n"
         "✓  영구 저장 책임은 클라이언트(IndexedDB), 서버는 분석 파이프라인만 책임",
         size=13, color=NAVY)

footer(s, 9, TOTAL)

# =========================================================================
# Slide 10 — Why PostgreSQL
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "왜 PostgreSQL인가",
           subtitle="비디오는 안 들어갑니다. '의미 있는 텍스트/JSON만 저장' 영역에 정확히 맞아떨어진 선택.",
           kicker="WHY POSTGRES")

# left: what we store
add_round(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(4.5),
          WHITE, line=GREY_LIGHT)
add_text(s, Inches(0.9), Inches(2.25), Inches(5.5), Inches(0.4),
         "STORE  ✓", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(2.65), Inches(5.5), Inches(3.7),
         "•  users · projects · sessions (관계형)\n"
         "•  agent_messages — LLM 대화 로그 (JSONB)\n"
         "•  종합 리포트 — axis 점수 · annotated moments (JSONB)\n"
         "•  세션 메타 · 시나리오 · focus_goals",
         size=13, color=GREY_DARK)

add_text(s, Inches(0.9), Inches(5.6), Inches(5.5), Inches(0.4),
         "NOT STORE  ✗", size=12, bold=True, color=RED)
add_text(s, Inches(0.9), Inches(5.95), Inches(5.5), Inches(0.5),
         "영상 blob — IndexedDB가 담당",
         size=13, bold=True, color=GREY_DARK)

# right: 3 reasons
reasons = [
    ("JSONB", "리포트·LLM 응답·이벤트 같은\n스키마 유연 데이터를 그대로\n인덱싱·쿼리"),
    ("단일 인스턴스", "관계형 + 문서형 동시 해결\nMongo + RDB 이중화 불필요"),
    ("Docker 친화", "Compose에 그대로 얹기 좋음\n무료 라이선스 · 운영 단순성"),
]
for i, (t, d) in enumerate(reasons):
    cy = Inches(2.1) + Inches(i * 1.5)
    add_round(s, Inches(6.8), cy, Inches(5.9), Inches(1.35), ACCENT_SOFT)
    add_rect(s, Inches(6.8), cy, Inches(0.16), Inches(1.35), ACCENT)
    add_text(s, Inches(7.05), cy + Inches(0.15), Inches(5.5), Inches(0.4),
             t, size=16, bold=True, color=NAVY)
    add_text(s, Inches(7.05), cy + Inches(0.55), Inches(5.5), Inches(0.8),
             d, size=11, color=GREY_DARK)

footer(s, 10, TOTAL)

# =========================================================================
# Slide 11 — 시나리오 확장성
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "시나리오 확장성 — YAML 한 장 = 새 코칭",
           subtitle="services/coach/rubrics/*.yaml 만 추가하면 코드 변경 0줄로 새 시나리오 동작",
           kicker="EXTENSIBILITY")

# left: 3 yaml chips
yamls = [
    ("presentation.yaml", "logic · delivery · gaze\n자세 · 표정", ACCENT),
    ("interview.yaml", "자신감 · 답변 명료성\n시선 · 페이스", HIGHLIGHT),
    ("vocal.yaml", "피치 다양성 · 발음\n리듬 · 호흡 · 표정", RED),
]
for i, (n, axes, col) in enumerate(yamls):
    cy = Inches(2.1) + Inches(i * 1.5)
    add_round(s, Inches(0.6), cy, Inches(5.9), Inches(1.35), WHITE, line=GREY_LIGHT)
    add_rect(s, Inches(0.6), cy, Inches(0.16), Inches(1.35), col)
    add_text(s, Inches(0.85), cy + Inches(0.15), Inches(5.5), Inches(0.45),
             n, size=15, bold=True, color=NAVY, font=FONT_MONO)
    add_text(s, Inches(0.85), cy + Inches(0.6), Inches(5.5), Inches(0.7),
             axes, size=12, color=GREY_DARK)

# right: code snippet style
add_round(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.5),
          NAVY, radius=0.05)
add_text(s, Inches(7.05), Inches(2.25), Inches(5.5), Inches(0.4),
         "rubrics/presentation.yaml (발췌)", size=11, bold=True, color=ACCENT)
yaml_text = (
    "id: presentation\n"
    "axes:\n"
    "  - id: logic\n"
    "    weight: 1.0\n"
    "    rubric_text: |\n"
    "      도입/본론/결론 명확성 …\n"
    "  - id: delivery\n"
    "    weight: 1.0\n"
    "event_kind_emphasis:\n"
    "  chin_on_hand: blunder\n"
    "  head_tilt_sustained: mistake\n"
    "top_priority_count: 3"
)
add_text(s, Inches(7.05), Inches(2.7), Inches(5.5), Inches(3.8),
         yaml_text, size=11, color=ACCENT_SOFT, font=FONT_MONO)

footer(s, 11, TOTAL)

# =========================================================================
# Slide 12 — Demo
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "Live Demo", kicker="SHOW, DON'T TELL")

# Demo flow steps
demo_steps = [
    ("①", "시나리오 선택", "URL ?scenario=presentation\n→ aggregator → coach"),
    ("②", "1분 녹화", "MediaPipe 5fps + Web Speech\n실시간 HUD 코멘트"),
    ("③", "정지 & 업로드", "IndexedDB → /analyze\nfaster-whisper + librosa"),
    ("④", "종합 리포트", "axis 점수 · annotated moments\n훈련 처방"),
]
x0 = Inches(0.6)
for i, (n, t, d) in enumerate(demo_steps):
    cx = x0 + Inches(i * 3.05)
    add_round(s, cx, Inches(2.3), Inches(2.95), Inches(3.3), WHITE, line=GREY_LIGHT)
    add_text(s, cx, Inches(2.5), Inches(2.95), Inches(0.6),
             n, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, cx + Inches(0.2), Inches(3.5), Inches(2.55), Inches(0.5),
             t, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, cx + Inches(0.2), Inches(4.1), Inches(2.55), Inches(1.4),
             d, size=11, color=GREY_MID, align=PP_ALIGN.CENTER)

add_round(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.7),
          NAVY, radius=0.3)
add_text(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.7),
         "Backup : 캡처 3장 (라이브 HUD / 리포트 / 아바타 재생) 준비",
         size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 12, TOTAL)

# =========================================================================
# Slide 13 — 한계 & V2
# =========================================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, "한계와 V2 로드맵",
           subtitle="V1은 의도적으로 좁게 잘랐습니다. 그 경계가 곧 다음 작업 목록.",
           kicker="LIMITS & NEXT")

# left V1 빠진 것
add_round(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(4.5),
          GREY_LIGHT)
add_text(s, Inches(0.9), Inches(2.25), Inches(5.5), Inches(0.4),
         "V1에서 의도적으로 뺀 것", size=13, bold=True, color=GREY_MID)
add_text(s, Inches(0.9), Inches(2.7), Inches(5.5), Inches(3.8),
         "•  청크 / Resumable 업로드\n"
         "•  다중 사용자 협업\n"
         "•  아바타 UI 노출 (코드는 유지)\n"
         "•  SER (감정 분류 모델)\n"
         "•  WhisperX forced alignment\n"
         "•  향상도 비교 / 개인화 추적",
         size=14, color=GREY_DARK, line_spacing=1.5)

# right V2 후보
add_round(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.5),
          ACCENT_SOFT)
add_text(s, Inches(7.1), Inches(2.25), Inches(5.5), Inches(0.4),
         "V2 후보", size=13, bold=True, color=ACCENT)
add_text(s, Inches(7.1), Inches(2.7), Inches(5.5), Inches(3.8),
         "•  pyannote — 다중 화자 (대화 시나리오)\n"
         "•  OpenFace — 정밀 Action Unit · gaze\n"
         "•  wav2vec2 SER — 톤/감정 분류\n"
         "•  Resumable upload (tus.io)\n"
         "•  아바타 익명 모드 토글\n"
         "•  세션 간 향상도 대시보드",
         size=14, bold=True, color=NAVY, line_spacing=1.5)

footer(s, 13, TOTAL)

# =========================================================================
# Slide 14 — Q&A
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
c.fill.transparency = 0.6

add_text(s, Inches(0.9), Inches(2.3), Inches(12), Inches(1.5),
         "Q & A", size=96, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(4.1), Inches(12), Inches(0.6),
         "Thank you", size=24, color=ACCENT_SOFT)
add_text(s, Inches(0.9), Inches(4.7), Inches(12), Inches(0.5),
         "백업 슬라이드: 백엔드 토폴로지 · WebSocket JSON · 트리거 cooldown 표 · rubric YAML 스니펫",
         size=12, color=ACCENT_SOFT)

out = r"c:\Users\swh01\presentation-coach\SpeakUp_Capstone_TechStack.pptx"
prs.save(out)
print(f"Saved: {out}")
